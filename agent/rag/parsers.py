from __future__ import annotations

from pathlib import Path


def _read_text_file(path: Path) -> str:
    try:
        return path.read_text(encoding="utf-8")
    except UnicodeDecodeError:
        return path.read_text(encoding="gbk", errors="ignore")


def parse_txt(path: Path) -> str:
    return _read_text_file(path)


def parse_pdf(path: Path) -> str:
    """解析 PDF 文件，提取文本内容"""
    try:
        from pypdf import PdfReader
    except ImportError as exc:
        raise RuntimeError("pypdf not installed. Install `pypdf`.") from exc
    reader = PdfReader(str(path))
    parts = []
    for page in reader.pages:
        text = page.extract_text() or ""
        if text:
            parts.append(text)
    return "\n".join(parts)


def extract_pdf_images(path: Path, output_dir: Path | None = None) -> list[dict]:
    """
    从 PDF 中提取所有嵌入的图片

    Args:
        path: PDF 文件路径
        output_dir: 可选，保存图片的目录。如果为 None，返回 base64 编码的图片数据

    Returns:
        图片信息列表，每个元素包含:
        - page: 页码 (从1开始)
        - index: 图片在该页的索引
        - width: 图片宽度
        - height: 图片高度
        - path: 图片保存路径 (如果 output_dir 不为 None)
        - base64: base64 编码的图片数据 (如果 output_dir 为 None)
        - media_type: 图片 MIME 类型
    """
    try:
        import fitz  # PyMuPDF
    except ImportError as exc:
        raise RuntimeError("PyMuPDF not installed. Install `pymupdf`.") from exc

    import base64

    doc = fitz.open(str(path))
    images = []

    for page_num, page in enumerate(doc, 1):
        image_list = page.get_images(full=True)

        for img_idx, img_info in enumerate(image_list):
            xref = img_info[0]  # 图片的 xref (交叉引用号)

            try:
                base_image = doc.extract_image(xref)
                if not base_image:
                    continue

                image_bytes = base_image["image"]
                image_ext = base_image["ext"]  # 图片扩展名 (png, jpeg, etc.)
                width = base_image.get("width", 0)
                height = base_image.get("height", 0)

                # 确定 MIME 类型
                mime_map = {
                    "png": "image/png",
                    "jpeg": "image/jpeg",
                    "jpg": "image/jpeg",
                    "gif": "image/gif",
                    "bmp": "image/bmp",
                    "tiff": "image/tiff",
                    "webp": "image/webp",
                }
                media_type = mime_map.get(image_ext.lower(), f"image/{image_ext}")

                img_data = {
                    "page": page_num,
                    "index": img_idx,
                    "width": width,
                    "height": height,
                    "media_type": media_type,
                }

                if output_dir:
                    # 保存到文件
                    output_dir.mkdir(parents=True, exist_ok=True)
                    filename = f"page{page_num}_img{img_idx}.{image_ext}"
                    img_path = output_dir / filename
                    with open(img_path, "wb") as f:
                        f.write(image_bytes)
                    img_data["path"] = str(img_path)
                else:
                    # 返回 base64 编码
                    img_data["base64"] = base64.standard_b64encode(image_bytes).decode("utf-8")

                images.append(img_data)

            except Exception:
                # 跳过无法提取的图片
                continue

    doc.close()
    return images


def parse_pdf_with_images(path: Path) -> dict:
    """
    解析 PDF 文件，同时提取文本和图片

    Returns:
        包含 text 和 images 的字典:
        - text: 提取的文本内容
        - images: 图片列表 (base64 编码)
    """
    text = parse_pdf(path)
    images = extract_pdf_images(path, output_dir=None)
    return {
        "text": text,
        "images": images,
    }


def render_pdf_pages(path: Path, page_numbers: list[int] = None, dpi: int = 150) -> list[dict]:
    """
    将 PDF 页面渲染为图片（整页截图）

    与 extract_pdf_images 不同，这个函数把整个页面渲染成图片，
    包括矢量图形、文本等所有内容。适用于查看 PDF 页面布局。

    Args:
        path: PDF 文件路径
        page_numbers: 要渲染的页码列表 (从1开始)，None 表示所有页面
        dpi: 渲染分辨率，默认 150

    Returns:
        渲染后的图片列表，每个元素包含:
        - page: 页码
        - width: 图片宽度
        - height: 图片高度
        - base64: base64 编码的 PNG 图片
        - media_type: "image/png"
    """
    try:
        import fitz  # PyMuPDF
    except ImportError as exc:
        raise RuntimeError("PyMuPDF not installed. Install `pymupdf`.") from exc

    import base64

    doc = fitz.open(str(path))
    total_pages = len(doc)
    results = []

    # 确定要渲染的页面
    if page_numbers is None:
        pages_to_render = range(1, total_pages + 1)
    else:
        pages_to_render = [p for p in page_numbers if 1 <= p <= total_pages]

    # 计算缩放因子 (72 是 PDF 默认 DPI)
    zoom = dpi / 72.0
    matrix = fitz.Matrix(zoom, zoom)

    for page_num in pages_to_render:
        page = doc[page_num - 1]  # fitz 使用 0-based 索引

        # 渲染页面为 PNG
        pix = page.get_pixmap(matrix=matrix)
        png_data = pix.tobytes("png")

        results.append({
            "page": page_num,
            "width": pix.width,
            "height": pix.height,
            "base64": base64.standard_b64encode(png_data).decode("utf-8"),
            "media_type": "image/png",
        })

    doc.close()
    return results


def get_pdf_info(path: Path) -> dict:
    """
    获取 PDF 文件的基本信息

    Returns:
        - total_pages: 总页数
        - title: 标题 (如果有)
        - author: 作者 (如果有)
    """
    try:
        import fitz
    except ImportError as exc:
        raise RuntimeError("PyMuPDF not installed. Install `pymupdf`.") from exc

    doc = fitz.open(str(path))

    info = {
        "total_pages": len(doc),
        "title": doc.metadata.get("title", ""),
        "author": doc.metadata.get("author", ""),
    }

    doc.close()
    return info


def parse_docx(path: Path) -> str:
    try:
        import docx
    except ImportError as exc:
        raise RuntimeError("python-docx not installed. Install `python-docx`.") from exc
    document = docx.Document(str(path))
    return "\n".join(p.text for p in document.paragraphs if p.text)


def parse_xlsx(path: Path) -> str:
    try:
        import openpyxl
    except ImportError as exc:
        raise RuntimeError("openpyxl not installed. Install `openpyxl`.") from exc
    workbook = openpyxl.load_workbook(path, data_only=True, read_only=True)
    lines = []
    for sheet in workbook.worksheets:
        lines.append(f"[sheet:{sheet.title}]")
        for row in sheet.iter_rows(values_only=True):
            values = ["" if cell is None else str(cell) for cell in row]
            line = "\t".join(values).strip()
            if line:
                lines.append(line)
    return "\n".join(lines)


def parse_pptx(path: Path) -> str:
    """解析 PPTX 文件，提取所有文本内容"""
    try:
        from pptx import Presentation
    except ImportError as exc:
        raise RuntimeError("python-pptx not installed. Install `python-pptx`.") from exc

    prs = Presentation(str(path))
    lines = []

    for slide_num, slide in enumerate(prs.slides, 1):
        slide_texts = []

        for shape in slide.shapes:
            # 处理文本框
            if shape.has_text_frame:
                for paragraph in shape.text_frame.paragraphs:
                    text = paragraph.text.strip()
                    if text:
                        slide_texts.append(text)

            # 处理表格
            if shape.has_table:
                table = shape.table
                for row in table.rows:
                    row_texts = []
                    for cell in row.cells:
                        cell_text = cell.text.strip()
                        if cell_text:
                            row_texts.append(cell_text)
                    if row_texts:
                        slide_texts.append("\t".join(row_texts))

        if slide_texts:
            lines.append(f"[slide:{slide_num}]")
            lines.extend(slide_texts)

    return "\n".join(lines)


_PLAIN_TEXT_SUFFIXES = {
    ".txt", ".md", ".markdown", ".rst",
    ".py", ".js", ".ts", ".jsx", ".tsx", ".java", ".go", ".rs", ".c", ".h",
    ".cpp", ".hpp", ".cs", ".rb", ".php", ".sh", ".bat", ".ps1", ".sql",
    ".json", ".yaml", ".yml", ".toml", ".ini", ".cfg", ".conf", ".env",
    ".html", ".htm", ".xml", ".css", ".scss", ".less",
    ".csv", ".tsv", ".log",
}


def extract_text(path: Path) -> str:
    """提取文件文本内容。

    Binary formats: pdf, docx, xlsx, pptx.
    Plain-text formats (see _PLAIN_TEXT_SUFFIXES): read as UTF-8 (GBK fallback).
    """
    suffix = path.suffix.lower()
    if suffix == ".pdf":
        return parse_pdf(path)
    if suffix == ".docx":
        return parse_docx(path)
    if suffix == ".xlsx":
        return parse_xlsx(path)
    if suffix == ".pptx":
        return parse_pptx(path)
    if suffix in _PLAIN_TEXT_SUFFIXES:
        return _read_text_file(path)
    raise ValueError(f"Unsupported file type: {path}")
