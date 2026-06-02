"""P12.2.3 PowerPoint live dual-model behavior smoke.

Drives both doubao-code and gpt-5.5 through the same PPT-creation prompt
using the real `PowerPointRuntimeEdit` tool and a real PowerPoint COM
backend. Captures the full tool-call sequence (which ops, in what order,
in what batches), then probes the produced .pptx with python-pptx to
record the actual delivered structure.

The runner auto-accepts every diff preview (the trust loop wiring is
already covered by the P12.4.x UI smoke; here we want to observe what
the two models *do* under identical conditions, not re-test the UI).

Outputs:
  tests/results/p12_pptx_dual_model_live/<ts>/
    doubao-code/{raw_sse.txt, events.jsonl, assistant.txt,
                 summary.json, output.pptx, pptx_structure.json}
    gpt-5.5/{...same...}
    comparison_report.md
    summary.json
"""

from __future__ import annotations

import argparse
import json
import os
import re
import shutil
import socket
import subprocess
import sys
import threading
import time
import urllib.request
from collections import Counter
from contextlib import closing
from pathlib import Path
from typing import Any

import requests

ROOT = Path(__file__).resolve().parents[2]
if str(ROOT) not in sys.path:
    sys.path.insert(0, str(ROOT))

RESULTS = ROOT / "tests" / "results" / "p12_pptx_dual_model_live"
APP_YAML = ROOT / "config" / "app.yaml"


PROMPT_TEMPLATE = (
    "请帮我做一份 PowerPoint pptx 幻灯片演示文稿，路径已经定好："
    "{path}\n"
    "整个 deck 一共 3 slides：\n"
    "  Slide 1：封面 slide —— 居中文本框写「Quarterly Review 2026 Q2」。\n"
    "  Slide 2：要点 slide —— 一个文本框，里面写三行：\n"
    "    - Revenue +18% YoY\n"
    "    - New customers +540\n"
    "    - Churn 2.1%\n"
    "  Slide 3：流程图 slide —— 画一个矩形 + 一个圆形，分别写「Plan」和「Ship」，"
    "中间用直线 connector 连起来。\n"
    "完成后调用 save。文件路径已经给你，不需要先 Glob 或 Read 任何东西。"
    "请全程使用 PowerPointRuntimeEdit 这一个工具；可以分多次调用，也可以批一次性提交 ops。"
)


def _ts() -> str:
    return time.strftime("%Y%m%d_%H%M%S")


def _free_port() -> int:
    with closing(socket.socket(socket.AF_INET, socket.SOCK_STREAM)) as sock:
        sock.bind(("127.0.0.1", 0))
        return int(sock.getsockname()[1])


def _wait_for_server(base_url: str, timeout_s: float = 30.0) -> None:
    deadline = time.time() + timeout_s
    last = ""
    while time.time() < deadline:
        try:
            with urllib.request.urlopen(
                f"{base_url}/api/agent_runtime", timeout=3
            ) as resp:
                if 200 <= resp.status < 300:
                    return
        except Exception as exc:
            last = str(exc)
        time.sleep(0.5)
    raise RuntimeError(f"server did not start: {last}")


def _swap_profile(target: str) -> str:
    text = APP_YAML.read_text(encoding="utf-8")
    lines = text.splitlines()
    old = ""
    for i, line in enumerate(lines):
        if line.startswith("active_profile:"):
            old = line.split(":", 1)[1].strip()
            lines[i] = f"active_profile: {target}"
            break
    APP_YAML.write_text("\n".join(lines) + "\n", encoding="utf-8")
    return old


def _parse_sse(blob: str) -> list[dict]:
    out: list[dict] = []
    event = ""
    data_lines: list[str] = []

    def flush() -> None:
        nonlocal event, data_lines
        if not data_lines:
            event = ""
            return
        try:
            payload = json.loads("\n".join(data_lines))
        except Exception:
            payload = {"raw": "\n".join(data_lines)}
        out.append({"event": event or "message", "data": payload})
        event = ""
        data_lines = []

    for line in blob.splitlines():
        if line == "":
            flush()
            continue
        if line.startswith("event:"):
            event = line[6:].strip()
        elif line.startswith("data:"):
            data_lines.append(line[5:].strip())
    flush()
    return out


def _create_conversation(base_url: str, title: str) -> str:
    res = requests.post(
        f"{base_url}/api/conversations", json={"title": title}, timeout=10
    )
    res.raise_for_status()
    body = res.json()
    return str(body.get("conversation_id") or body.get("id"))


def _start_server(out_dir: Path, port: int) -> subprocess.Popen:
    return subprocess.Popen(
        [
            sys.executable, "-m", "uvicorn", "--factory",
            "--host", "127.0.0.1", "--port", str(port),
            "--app-dir", str(ROOT), "agent.ui.server:create_app",
        ],
        env=dict(os.environ),
        cwd=str(ROOT),
        stdout=(out_dir / "server_stdout.txt").open(
            "a", encoding="utf-8", errors="replace"
        ),
        stderr=(out_dir / "server_stderr.txt").open(
            "a", encoding="utf-8", errors="replace"
        ),
        text=True,
    )


def _probe_pptx(path: Path) -> dict[str, Any]:
    """Return a structured snapshot of the produced .pptx, or an error blob."""
    if not path.exists():
        return {"exists": False, "error": "file not found"}
    try:
        from pptx import Presentation
        from pptx.util import Emu
    except Exception as exc:
        return {"exists": True, "error": f"pptx import failed: {exc}"}
    try:
        prs = Presentation(str(path))
    except Exception as exc:
        return {"exists": True, "error": f"open failed: {exc}"}
    slides: list[dict[str, Any]] = []
    shape_total = 0
    has_rectangle = False
    has_oval = False
    has_connector = False
    for s_idx, slide in enumerate(prs.slides, start=1):
        shapes: list[dict[str, Any]] = []
        for shape in slide.shapes:
            shape_total += 1
            stype = str(shape.shape_type) if shape.shape_type is not None else ""
            auto_type = ""
            try:
                auto_raw = getattr(shape, "auto_shape_type", None)
                if auto_raw is not None:
                    auto_type = str(auto_raw)
            except Exception:
                auto_type = ""
            text = ""
            try:
                if shape.has_text_frame:
                    text = (shape.text_frame.text or "").strip()
            except Exception:
                pass
            try:
                left = int(shape.left) if shape.left is not None else None
                top = int(shape.top) if shape.top is not None else None
                width = int(shape.width) if shape.width is not None else None
                height = int(shape.height) if shape.height is not None else None
            except Exception:
                left = top = width = height = None
            name = getattr(shape, "name", "") or ""
            shapes.append({
                "shape_type": stype,
                "auto_shape_type": auto_type,
                "name": name,
                "text_preview": text[:120],
                "bbox_emu": [left, top, width, height],
            })
            # Recognise rectangle/oval/connector across three signals:
            #   - shape_type string (works for LINE / CONNECTOR)
            #   - auto_shape_type (works for AUTO_SHAPE rectangle vs oval)
            #   - shape name as fallback (model often names them semantically)
            haystack = " ".join([stype, auto_type, name]).lower()
            if "rectangle" in haystack or "rect" in haystack:
                has_rectangle = True
            if "oval" in haystack or "ellipse" in haystack or "circle" in haystack:
                has_oval = True
            if "line" in haystack or "connector" in haystack:
                has_connector = True
        slides.append({"index": s_idx, "shape_count": len(shapes), "shapes": shapes})
    return {
        "exists": True,
        "slide_count": len(slides),
        "shape_total": shape_total,
        "has_rectangle": has_rectangle,
        "has_oval_or_circle": has_oval,
        "has_connector_or_line": has_connector,
        "slides": slides,
    }


def _run_case(
    *, profile: str, base_url: str, out_dir: Path, prompt: str, target_pptx: Path
) -> dict[str, Any]:
    conv_id = _create_conversation(base_url, f"P12.2.3 {profile}")
    payload = {
        "message": prompt,
        "mode": "restricted",
        "plan_mode": False,
        "conversation_id": conv_id,
        "history": [],
        "max_iterations": 12,
    }

    raw_chunks: list[str] = []
    approved_ids: set[str] = set()
    _id_re = re.compile(r'"preview_id"\s*:\s*"([a-zA-Z0-9_-]+)"')

    def _maybe_approve(chunk: str) -> None:
        for m in _id_re.finditer(chunk):
            pid = m.group(1)
            if pid in approved_ids:
                continue
            approved_ids.add(pid)
            def _post() -> None:
                try:
                    requests.post(
                        f"{base_url}/api/diff_previews/{pid}",
                        json={"approved": True, "note": "auto-approve (P12.2.3 live)"},
                        timeout=20,
                    )
                except Exception:
                    pass
            threading.Thread(target=_post, daemon=True).start()

    start = time.time()
    error_msg = ""
    try:
        with requests.post(
            f"{base_url}/api/agent_chat_v2",
            json=payload, stream=True, timeout=600,
        ) as resp:
            resp.raise_for_status()
            for chunk in resp.iter_content(chunk_size=None, decode_unicode=True):
                if not chunk:
                    continue
                if isinstance(chunk, bytes):
                    chunk = chunk.decode("utf-8", errors="replace")
                raw_chunks.append(chunk)
                _maybe_approve(chunk)
    except Exception as exc:
        error_msg = f"{type(exc).__name__}: {exc}"
    elapsed = round(time.time() - start, 2)

    blob = "".join(raw_chunks)
    (out_dir / "raw_sse.txt").write_text(blob, encoding="utf-8")
    events = _parse_sse(blob)
    (out_dir / "events.jsonl").write_text(
        "\n".join(json.dumps(e, ensure_ascii=False) for e in events),
        encoding="utf-8",
    )
    activities = [e["data"] for e in events if e["event"] == "activity"]
    tokens = [e["data"].get("text", "") for e in events if e["event"] == "token"]
    done = next((e["data"] for e in events if e["event"] == "done"), {})
    assistant = "".join(tokens)
    (out_dir / "assistant.txt").write_text(assistant, encoding="utf-8")

    # Build tool-call timeline by inspecting activity events.
    timeline: list[dict[str, Any]] = []
    pptx_calls: list[dict[str, Any]] = []
    other_tool_calls: list[str] = []
    diff_preview_count = 0
    for act in activities:
        atype = act.get("type")
        meta = act.get("meta") or {}
        if atype == "tool_call":
            name = meta.get("name") or ""
            args_preview = meta.get("args_preview") or meta.get("input_preview") or ""
            timeline.append({"event": "tool_call", "name": name, "args_preview": args_preview})
            if name == "PowerPointRuntimeEdit":
                ops = []
                # Args preview may be a JSON-like string or the structured input.
                inp = meta.get("input") if isinstance(meta.get("input"), dict) else None
                if inp and isinstance(inp.get("ops"), list):
                    ops = [str(o.get("op")) for o in inp["ops"] if isinstance(o, dict)]
                else:
                    # Fall back to scanning args_preview text for op kinds.
                    ops = re.findall(r'"op"\s*:\s*"([a-z_]+)"', str(args_preview))
                pptx_calls.append({"ops": ops, "op_count": len(ops)})
            else:
                other_tool_calls.append(name)
        elif atype == "diff_preview":
            diff_preview_count += 1
            timeline.append({
                "event": "diff_preview",
                "tool": meta.get("tool"),
                "op_count": meta.get("op_count"),
            })
        elif atype == "tool_result":
            timeline.append({
                "event": "tool_result",
                "name": (meta.get("name") or ""),
                "is_error": bool(meta.get("is_error")),
            })

    op_kind_counter: Counter[str] = Counter()
    for call in pptx_calls:
        op_kind_counter.update(call["ops"])
    op_counts_per_call = [c["op_count"] for c in pptx_calls]

    pptx_struct = _probe_pptx(target_pptx)
    (out_dir / "pptx_structure.json").write_text(
        json.dumps(pptx_struct, ensure_ascii=False, indent=2),
        encoding="utf-8",
    )

    task_complete = (
        pptx_struct.get("exists") is True
        and pptx_struct.get("slide_count", 0) >= 3
        and pptx_struct.get("has_rectangle") is True
    )

    summary = {
        "profile": profile,
        "elapsed_s": elapsed,
        "error": error_msg,
        "assistant_text_length": len(assistant),
        "done": done,
        "diff_preview_count": diff_preview_count,
        "pptx_call_count": len(pptx_calls),
        "max_ops_in_one_call": max(op_counts_per_call) if op_counts_per_call else 0,
        "total_ops": sum(op_counts_per_call),
        "op_kind_distribution": dict(op_kind_counter),
        "other_tool_calls": sorted(set(other_tool_calls)),
        "timeline": timeline,
        "task_complete": task_complete,
        "pptx_structure": pptx_struct,
    }
    (out_dir / "summary.json").write_text(
        json.dumps(summary, ensure_ascii=False, indent=2), encoding="utf-8"
    )
    return summary


def _build_report(out_root: Path, results: dict[str, dict[str, Any]]) -> str:
    def cell(v: Any) -> str:
        if isinstance(v, dict):
            return ", ".join(f"{k}={v[k]}" for k in sorted(v))
        if isinstance(v, list):
            return ", ".join(map(str, v)) if v else "—"
        if v is None:
            return "—"
        return str(v)

    profiles = list(results.keys())
    if not profiles:
        return "# No results.\n"

    rows = [
        ("Profile", lambda r: r.get("profile")),
        ("Elapsed (s)", lambda r: r.get("elapsed_s")),
        ("Task complete (≥3 slides + rectangle)", lambda r: r.get("task_complete")),
        ("PowerPointRuntimeEdit calls", lambda r: r.get("pptx_call_count")),
        ("Total ops emitted", lambda r: r.get("total_ops")),
        ("Max ops in one call", lambda r: r.get("max_ops_in_one_call")),
        ("Diff previews shown", lambda r: r.get("diff_preview_count")),
        ("op_kind distribution", lambda r: r.get("op_kind_distribution")),
        ("Other tools called", lambda r: r.get("other_tool_calls")),
        ("Slide count (in file)", lambda r: (r.get("pptx_structure") or {}).get("slide_count")),
        ("Shape total (in file)", lambda r: (r.get("pptx_structure") or {}).get("shape_total")),
        ("Has rectangle", lambda r: (r.get("pptx_structure") or {}).get("has_rectangle")),
        ("Has oval/circle", lambda r: (r.get("pptx_structure") or {}).get("has_oval_or_circle")),
        ("Has connector/line", lambda r: (r.get("pptx_structure") or {}).get("has_connector_or_line")),
        ("Error", lambda r: r.get("error") or "—"),
    ]

    lines = ["# P12.2.3 — PowerPoint Live Dual-Model Behavior\n"]
    lines.append(f"Run dir: `{out_root}`\n")
    lines.append("")
    header = "| Metric | " + " | ".join(profiles) + " |"
    sep = "|" + "---|" * (len(profiles) + 1)
    lines.append(header)
    lines.append(sep)
    for label, fn in rows:
        cells = [cell(fn(results[p])) for p in profiles]
        lines.append(f"| {label} | " + " | ".join(cells) + " |")

    lines.append("\n## Tool-call timeline (first 25 events per profile)\n")
    for p in profiles:
        lines.append(f"### {p}\n")
        timeline = results[p].get("timeline") or []
        for i, ev in enumerate(timeline[:25]):
            lines.append(f"- [{i}] **{ev.get('event')}** — {json.dumps({k: v for k, v in ev.items() if k != 'event'}, ensure_ascii=False)}")
        if len(timeline) > 25:
            lines.append(f"- … (+{len(timeline) - 25} more)")
        lines.append("")

    lines.append("## Behavior observations\n")
    if len(profiles) == 2:
        a, b = profiles[0], profiles[1]
        ra, rb = results[a], results[b]
        notes: list[str] = []

        def _safe(rec: dict, *path):
            cur: Any = rec
            for k in path:
                if not isinstance(cur, dict):
                    return None
                cur = cur.get(k)
            return cur

        ca = ra.get("pptx_call_count") or 0
        cb = rb.get("pptx_call_count") or 0
        if ca and cb:
            if ca == cb:
                notes.append(f"- 两个模型都用了 **{ca}** 次 `PowerPointRuntimeEdit`。")
            elif ca < cb:
                notes.append(
                    f"- **{a}** 用 {ca} 次（更倾向 batch），**{b}** 用 {cb} 次（更 incremental）。"
                )
            else:
                notes.append(
                    f"- **{b}** 用 {cb} 次（更倾向 batch），**{a}** 用 {ca} 次（更 incremental）。"
                )

        ma = ra.get("max_ops_in_one_call") or 0
        mb = rb.get("max_ops_in_one_call") or 0
        notes.append(
            f"- 单次最多 ops：**{a}** = {ma}，**{b}** = {mb}。差值 {abs(ma - mb)} 表示批处理偏好不同。"
        )

        oa = ra.get("op_kind_distribution") or {}
        ob = rb.get("op_kind_distribution") or {}
        only_a = sorted(set(oa) - set(ob))
        only_b = sorted(set(ob) - set(oa))
        if only_a:
            notes.append(f"- 只有 **{a}** 用到的 op：{', '.join(only_a)}")
        if only_b:
            notes.append(f"- 只有 **{b}** 用到的 op：{', '.join(only_b)}")
        if "save" in oa and "save" not in ob:
            notes.append(f"- ⚠️ **{b}** 整轮没调 `save`。")
        if "save" in ob and "save" not in oa:
            notes.append(f"- ⚠️ **{a}** 整轮没调 `save`。")
        if "create_presentation" in oa and "create_presentation" not in ob:
            notes.append(f"- **{a}** 主动调 `create_presentation`，**{b}** 假设文件存在。")
        if "create_presentation" in ob and "create_presentation" not in oa:
            notes.append(f"- **{b}** 主动调 `create_presentation`，**{a}** 假设文件存在。")

        if ra.get("task_complete") and not rb.get("task_complete"):
            notes.append(f"- ✅ **{a}** 任务完成；❌ **{b}** 没完成（看 pptx_structure.json）")
        elif rb.get("task_complete") and not ra.get("task_complete"):
            notes.append(f"- ✅ **{b}** 任务完成；❌ **{a}** 没完成（看 pptx_structure.json）")
        elif ra.get("task_complete") and rb.get("task_complete"):
            notes.append("- ✅ 两个模型都完成了任务（slide_count≥3 且产物含矩形）。")
        else:
            notes.append("- ❌ 两个模型都没完成任务。")

        ea = ra.get("elapsed_s") or 0
        eb = rb.get("elapsed_s") or 0
        if ea and eb:
            faster = a if ea < eb else b
            notes.append(
                f"- 用时：{a}={ea}s, {b}={eb}s（**{faster}** 更快，差 {abs(ea - eb):.1f}s）。"
            )

        lines.extend(notes)
    lines.append("")
    return "\n".join(lines)


def main() -> int:
    parser = argparse.ArgumentParser()
    parser.add_argument(
        "--profiles",
        default="doubao-code,gpt-5.5",
        help="Comma-separated profiles to drive (default: doubao-code,gpt-5.5)",
    )
    args = parser.parse_args()

    run_root = RESULTS / _ts()
    run_root.mkdir(parents=True, exist_ok=True)
    backup = run_root / "app.yaml.backup"
    shutil.copyfile(APP_YAML, backup)
    original: str | None = None

    summary: dict[str, Any] = {
        "run_dir": str(run_root),
        "results": {},
        "errors": [],
    }
    try:
        for profile in [p.strip() for p in args.profiles.split(",") if p.strip()]:
            case_dir = run_root / profile.replace("/", "_")
            case_dir.mkdir(parents=True, exist_ok=True)
            target_pptx = case_dir / "output.pptx"
            # Best-effort clean: PowerPoint will get a fresh path each time.
            if target_pptx.exists():
                try:
                    target_pptx.unlink()
                except Exception:
                    pass
            try:
                if original is None:
                    original = _swap_profile(profile)
                else:
                    _swap_profile(profile)
                port = _free_port()
                base_url = f"http://127.0.0.1:{port}"
                proc = _start_server(case_dir, port)
                try:
                    _wait_for_server(base_url)
                    case_summary = _run_case(
                        profile=profile,
                        base_url=base_url,
                        out_dir=case_dir,
                        prompt=PROMPT_TEMPLATE.format(path=str(target_pptx)),
                        target_pptx=target_pptx,
                    )
                    summary["results"][profile] = case_summary
                finally:
                    proc.terminate()
                    try:
                        proc.wait(timeout=10)
                    except subprocess.TimeoutExpired:
                        proc.kill()
            except Exception as exc:
                summary["errors"].append(f"{profile}: {type(exc).__name__}: {exc}")
                summary["results"][profile] = {
                    "profile": profile,
                    "error": str(exc),
                    "task_complete": False,
                }
    finally:
        if original:
            try:
                _swap_profile(original)
            except Exception as exc:
                summary["errors"].append(f"failed to restore profile: {exc}")
                shutil.copyfile(backup, APP_YAML)

        report_md = _build_report(run_root, summary["results"])
        (run_root / "comparison_report.md").write_text(report_md, encoding="utf-8")
        (run_root / "summary.json").write_text(
            json.dumps(summary, ensure_ascii=False, indent=2), encoding="utf-8"
        )

    print(json.dumps({
        "run_dir": summary["run_dir"],
        "errors": summary["errors"],
        "task_complete": {
            p: r.get("task_complete") for p, r in summary["results"].items()
        },
        "elapsed_s": {
            p: r.get("elapsed_s") for p, r in summary["results"].items()
        },
        "pptx_call_count": {
            p: r.get("pptx_call_count") for p, r in summary["results"].items()
        },
        "report_md": str(run_root / "comparison_report.md"),
    }, ensure_ascii=False, indent=2))
    all_complete = all(
        r.get("task_complete") for r in summary["results"].values()
    ) and not summary["errors"]
    return 0 if all_complete else 1


if __name__ == "__main__":
    raise SystemExit(main())
