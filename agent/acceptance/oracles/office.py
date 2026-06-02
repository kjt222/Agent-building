"""Office L2 oracle (P14.2.3): thin docx / xlsx / pptx parse + probes.

The Office runtime tools (`agent/tools_v2/word_tool.py` etc.) already
maintain strong structural invariants on write. This oracle's job is to
catch the **post-hoc** divergence: agent claimed it added a section /
sheet / slide, but parsing the saved file shows nothing was committed.

Checks per file type:
  - .docx: python-docx parse → paragraph count, table count, hyperlink count
  - .xlsx: openpyxl parse → sheet names, total non-empty cells
  - .pptx: python-pptx parse → slide count, total shape count

Failure mode is "file present but empty / unparseable" → fail. Findings
include the parsed key fields so the verdict.json carries forensic context
without re-opening the file.
"""

from __future__ import annotations

from pathlib import Path
from typing import Any, Iterable

from ..oracle import OracleReport, register_oracle


def _probe_docx(p: Path) -> tuple[bool, str, dict]:
    try:
        from docx import Document  # type: ignore
    except Exception as exc:
        return True, f"python-docx unavailable: {exc}", {"skipped": True}
    try:
        doc = Document(str(p))
    except Exception as exc:
        return False, f"docx parse failed: {exc}", {}
    paras = list(doc.paragraphs)
    return True, "", {
        "paragraph_count": len(paras),
        "table_count": len(doc.tables),
        "section_count": len(doc.sections),
        "non_empty_paragraph_count": sum(1 for p in paras if p.text.strip()),
    }


def _probe_xlsx(p: Path) -> tuple[bool, str, dict]:
    try:
        from openpyxl import load_workbook  # type: ignore
    except Exception as exc:
        return True, f"openpyxl unavailable: {exc}", {"skipped": True}
    try:
        wb = load_workbook(str(p), read_only=True, data_only=False)
    except Exception as exc:
        return False, f"xlsx parse failed: {exc}", {}
    sheet_names = list(wb.sheetnames)
    non_empty = 0
    for s in sheet_names:
        ws = wb[s]
        for row in ws.iter_rows(values_only=True):
            for c in row:
                if c is not None and (not isinstance(c, str) or c.strip()):
                    non_empty += 1
    return True, "", {
        "sheet_count": len(sheet_names),
        "sheet_names": sheet_names,
        "non_empty_cell_count": non_empty,
    }


def _probe_pptx(p: Path) -> tuple[bool, str, dict]:
    try:
        from pptx import Presentation  # type: ignore
    except Exception as exc:
        return True, f"python-pptx unavailable: {exc}", {"skipped": True}
    try:
        pres = Presentation(str(p))
    except Exception as exc:
        return False, f"pptx parse failed: {exc}", {}
    slide_count = len(pres.slides)
    shape_count = sum(len(slide.shapes) for slide in pres.slides)
    return True, "", {
        "slide_count": slide_count,
        "total_shape_count": shape_count,
    }


_PROBES = {
    ".docx": _probe_docx,
    ".xlsx": _probe_xlsx,
    ".pptx": _probe_pptx,
}


class OfficeOracle:
    name = "office"

    def check(
        self,
        artifact_paths: Iterable[Path],
        task_spec: dict[str, Any] | None = None,
    ) -> OracleReport:
        paths = [Path(p) for p in artifact_paths]
        if not paths:
            return OracleReport(
                oracle=self.name,
                verdict="unknown",
                findings=["no artifact paths provided"],
            )

        per_file: dict[str, dict] = {}
        findings: list[str] = []
        any_fail = False
        any_probed = False

        for p in paths:
            ext = p.suffix.lower()
            probe = _PROBES.get(ext)
            if probe is None:
                per_file[str(p)] = {"skipped": True, "reason": f"unsupported ext {ext!r}"}
                continue
            if not p.exists():
                per_file[str(p)] = {"missing": True}
                findings.append(f"[{p.name}] file missing")
                any_fail = True
                continue
            ok, err, key_fields = probe(p)
            per_file[str(p)] = {"ext": ext, **key_fields}
            if err and "unavailable" in err:
                # parser library not installed — skip rather than fail
                per_file[str(p)]["note"] = err
                continue
            any_probed = True
            if not ok:
                any_fail = True
                findings.append(f"[{p.name}] {err}")
                continue
            # Empty-content probe.
            if ext == ".docx" and key_fields.get("non_empty_paragraph_count", 0) == 0:
                findings.append(f"[{p.name}] docx has no non-empty paragraphs")
                any_fail = True
            elif ext == ".xlsx" and key_fields.get("non_empty_cell_count", 0) == 0:
                findings.append(f"[{p.name}] xlsx has no non-empty cells")
                any_fail = True
            elif ext == ".pptx" and key_fields.get("slide_count", 0) == 0:
                findings.append(f"[{p.name}] pptx has zero slides")
                any_fail = True

        if not any_probed and not any_fail:
            return OracleReport(
                oracle=self.name,
                verdict="unknown",
                findings=findings or ["no probable office file found"],
                evidence={"files": per_file},
            )

        verdict = "fail" if any_fail else "pass"
        return OracleReport(
            oracle=self.name,
            verdict=verdict,
            findings=findings,
            evidence={"files": per_file},
        )


register_oracle("office", OfficeOracle())
